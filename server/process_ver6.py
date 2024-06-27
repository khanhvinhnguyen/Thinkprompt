import sys
import os
import fitz 
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from reportlab.lib.utils import ImageReader
from reportlab.pdfbase.ttfonts import TTFont
from reportlab.pdfbase import pdfmetrics
from docx import Document
from docx.shared import Pt, RGBColor
from pptx import Presentation
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
from docx.oxml.ns import qn
from docx.oxml import OxmlElement
from translate import Translator

input_file = sys.argv[1]
output_file = sys.argv[2]

def ensure_dir_exists(directory):
    if not os.path.exists(directory):
        os.makedirs(directory)

def extract_pdf_text_images(input_file, output_dir):
    doc = fitz.open(input_file)
    full_text = ""
    images = []
    text_details = []

    for page_num in range(len(doc)):
        page = doc.load_page(page_num)
        blocks = page.get_text("dict")["blocks"]

        for block in blocks:
            if block["type"] == 0:
                for line in block["lines"]:
                    for span in line["spans"]:
                        text_details.append({
                            "text": span["text"],
                            "font": "Roboto-Medium",
                            "size": span["size"],
                            "color": span["color"],
                            "bbox": span["bbox"],
                            "page_num": page_num
                        })
                        full_text += span["text"] + "\n"
            elif block["type"] == 1: 
                img_list = page.get_images(full=True)
                for img_index, img in enumerate(img_list):
                    xref = img[0]
                    base_image = doc.extract_image(xref)
                    image_bytes = base_image["image"]
                    img_filename = f"image_{page_num+1}_{img_index+1}.png"
                    img_path = os.path.join(output_dir, "images", img_filename)
                    ensure_dir_exists(os.path.dirname(img_path))
                    with open(img_path, "wb") as img_file:
                        img_file.write(image_bytes)
                    images.append((img_path, block["bbox"], page_num)) 

    return full_text, images, text_details

def extract_docx_text_images(input_file, output_dir):
    doc = Document(input_file)
    full_text = ""
    images = []
    text_details = []
    tables = []

    for para in doc.paragraphs:
        alignment = para.alignment
        for run in para.runs:
            font = run.font
            text_details.append({
                "text": run.text,
                "font": "Roboto-Medium",
                "size": font.size.pt if font.size else None,
                "bold": font.bold,
                "italic": font.italic,
                "color": font.color.rgb if font.color else None,
                "alignment": alignment
            })
            full_text += run.text + "\n"

    for table in doc.tables:
        table_data = []
        for row in table.rows:
            row_data = []
            for cell in row.cells:
                cell_data = {
                    "text": cell.text,
                    "alignment": cell.paragraphs[0].alignment if cell.paragraphs else None
                }
                row_data.append(cell_data)
            table_data.append(row_data)
        tables.append(table_data)

    for rel in doc.part.rels:
        if "image" in rel:
            image = doc.part.rels[rel].target_part.blob
            img_filename = f"image_{len(images)+1}.png"
            img_path = os.path.join(output_dir, "images", img_filename)
            ensure_dir_exists(os.path.dirname(img_path))
            with open(img_path, "wb") as img_file:
                img_file.write(image)
            images.append(img_path)

    return full_text, images, text_details, tables

def convert_text_to_uppercase(full_text):
    return full_text.upper()

def create_uppercase_pdf(text_details, images, output_file):
    font_path = os.path.join(os.path.dirname(__file__), 'public', 'fonts', 'Roboto-Medium.ttf')
    pdfmetrics.registerFont(TTFont('Roboto-Medium', font_path))
    c = canvas.Canvas(output_file, pagesize=letter)
    width, height = letter
    margin = 50

    def add_text(c, text_details, max_width):
        x = margin 
        y = height - margin 
        for detail in text_details:
            c.setFont("Roboto-Medium", detail["size"])
            c.setFillColorRGB((detail["color"] >> 16 & 0xFF) / 255.0, (detail["color"] >> 8 & 0xFF) / 255.0, (detail["color"] & 0xFF) / 255.0)
            text = detail["text"].upper()
          
            if x + c.stringWidth(text, "Roboto-Medium", detail["size"]) > max_width:
                x = margin 
                y -= detail["size"] + 2 
            c.drawString(x, y, text)
            x += c.stringWidth(text + " ", "Roboto-Medium", detail["size"]) 

    add_text(c, text_details, width - 2 * margin)

    current_page = -1
    for img_path, bbox, page_num in images:
        if page_num != current_page:
            if current_page != -1:
                c.showPage()
            c.setPageSize(letter)
            c.setFont('Roboto-Medium', 12)
            current_page = page_num
        img = ImageReader(img_path)
        x0, y0, x1, y1 = bbox
        c.drawImage(img, x0, height - y1, width=(x1 - x0), height=(y1 - y0))

    c.save()
def set_cell_border(cell, border_color):
    for direction in ["top", "bottom", "start", "end"]:
        cell_border = OxmlElement(f'w:{direction}')
        cell_border.set(qn('w:val'), 'single')
        cell_border.set(qn('w:sz'), '4')
        cell_border.set(qn('w:space'), '0')
        cell_border.set(qn('w:color'), border_color)
        tcPr = cell._element.get_or_add_tcPr()
        tcPr.append(cell_border)

def create_uppercase_docx(text_details, images, tables, output_file):
    doc = Document()
    border_color = "0000FF" 

    for detail in text_details:
        para = doc.add_paragraph()
        para.alignment = detail["alignment"]
        run = para.add_run(detail["text"].upper())
        font = run.font
        font.name = "Roboto-Medium"
        if detail["size"]:
            font.size = Pt(detail["size"])
        if detail["bold"]:
            run.bold = True
        if detail["italic"]:
            run.italic = True
        if detail["color"]:
            font.color.rgb = detail["color"]

    for table_data in tables:
        table = doc.add_table(rows=len(table_data), cols=len(table_data[0]))
        for row_idx, row_data in enumerate(table_data):
            row = table.rows[row_idx]
            for col_idx, cell_data in enumerate(row_data):
                cell = row.cells[col_idx]
                cell.text = cell_data["text"].upper()
                if cell_data["alignment"] is not None:
                    cell.paragraphs[0].alignment = cell_data["alignment"]
                set_cell_border(cell, border_color)

    for img_path in images:
        doc.add_picture(img_path)

    doc.save(output_file)

def process_pptx(input_file, output_file):
    translator = Translator(to_lang="en")
    prs = Presentation(input_file)

    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    original_text = run.text
                    translated_text = translator.translate(original_text)
                    run.text = translated_text.upper()
                    run.font.size = Pt(12) if run.font.size is None else run.font.size
                    run.font.name = "Roboto-Medium"

    prs.save(output_file)

if input_file.lower().endswith('.pdf'):
    output_dir = os.path.dirname(output_file)
    ensure_dir_exists(output_dir)

    full_text, images, text_details = extract_pdf_text_images(input_file, output_dir)
    uppercase_text = convert_text_to_uppercase(full_text)
    create_uppercase_pdf(text_details, images, output_file)
elif input_file.lower().endswith('.docx'):
    output_dir = os.path.dirname(output_file)
    ensure_dir_exists(output_dir)

    full_text, images, text_details, tables = extract_docx_text_images(input_file, output_dir)
    create_uppercase_docx(text_details, images, tables, output_file)
elif input_file.lower().endswith('.pptx'):
    process_pptx(input_file, output_file)
else:
    print("Unsupported file format.")
